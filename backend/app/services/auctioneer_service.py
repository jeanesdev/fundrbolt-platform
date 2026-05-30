"""Service for auctioneer commission, settings, dashboard, and live auction logic."""

from __future__ import annotations

import logging
import re
from collections import defaultdict
from datetime import UTC, datetime
from decimal import Decimal
from html import unescape
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from app.schemas.revenue_generator import RevenueGeneratorDashboardSummary
from uuid import UUID

import httpx
from sqlalchemy import delete, func, select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.models.auction_bid import AuctionBid, BidStatus, PaddleRaiseContribution
from app.models.auction_item import AuctionItem, AuctionItemMedia
from app.models.auctioneer import AuctioneerEventSettings, AuctioneerItemCommission
from app.models.event import Event
from app.models.event_registration import EventRegistration
from app.models.quick_entry_bid import QuickEntryBid, QuickEntryBidStatus
from app.models.quick_entry_buy_now_bid import QuickEntryBuyNowBid
from app.models.quick_entry_donation import QuickEntryDonation
from app.models.quick_entry_donation_label import QuickEntryDonationLabelLink
from app.models.registration_guest import RegistrationGuest
from app.models.user import User
from app.schemas.auctioneer import (
    DEFAULT_PADDLE_RAISE_LEVELS,
    AuctioneerBidActivityEntry,
    AuctioneerBidderProfile,
    AuctioneerBidderSummary,
    AuctioneerItemDetailResponse,
    AuctioneerItemGalleryResponse,
    AuctioneerItemSummary,
    AuctioneerPaddleRaiseBidderSummary,
    AuctioneerPaddleRaiseLevelSummary,
    AuctioneerPaddleRaiseResponse,
    AuctionStatus,
    BidHistoryEntry,
    CommissionListItem,
    CommissionListResponse,
    CommissionResponse,
    DashboardResponse,
    EarningsSummary,
    EventSettingsResponse,
    EventTotals,
    HighBidder,
    LiveAuctionItem,
    LiveAuctionResponse,
    SilentAuctionStatus,
    TimerData,
)

logger = logging.getLogger(__name__)


class AuctioneerService:
    """Service for auctioneer-related operations."""

    def __init__(self, db: AsyncSession) -> None:
        self.db = db

    @staticmethod
    def _normalize_donation_label(label_name: str) -> str:
        return "Last Hero" if label_name.strip().lower() == "last leader" else label_name

    @staticmethod
    def _strip_slide_html(html: str | None) -> str:
        if not html:
            return ""

        normalized = re.sub(r"<br\s*/?>", "\n", html, flags=re.IGNORECASE)
        normalized = re.sub(r"</p\s*>", "\n\n", normalized, flags=re.IGNORECASE)
        normalized = re.sub(r"</li\s*>", "\n", normalized, flags=re.IGNORECASE)
        normalized = re.sub(r"<li[^>]*>", "• ", normalized, flags=re.IGNORECASE)
        normalized = re.sub(r"<[^>]+>", "", normalized)
        normalized = unescape(normalized)
        normalized = re.sub(r"\n{3,}", "\n\n", normalized)
        return normalized.strip()

    async def _get_primary_image_map(self, item_ids: list[UUID]) -> dict[UUID, str | None]:
        if not item_ids:
            return {}

        stmt = (
            select(AuctionItemMedia.auction_item_id, AuctionItemMedia.file_path)
            .where(
                AuctionItemMedia.auction_item_id.in_(item_ids),
                AuctionItemMedia.display_order == 0,
            )
            .order_by(AuctionItemMedia.display_order.asc(), AuctionItemMedia.created_at.asc())
        )
        rows = (await self.db.execute(stmt)).all()
        image_map: dict[UUID, str | None] = {}
        for item_id, file_path in rows:
            image_map.setdefault(item_id, file_path)
        return image_map

    async def _get_commission_map(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> dict[UUID, tuple[Decimal, Decimal]]:
        stmt = (
            select(
                AuctioneerItemCommission.auction_item_id,
                AuctioneerItemCommission.commission_percent,
                AuctioneerItemCommission.flat_fee,
            )
            .join(AuctionItem, AuctionItem.id == AuctioneerItemCommission.auction_item_id)
            .where(
                AuctionItem.event_id == event_id,
                AuctioneerItemCommission.auctioneer_user_id == auctioneer_user_id,
                AuctionItem.deleted_at.is_(None),
            )
        )
        return {
            auction_item_id: (Decimal(commission_percent), Decimal(flat_fee))
            for auction_item_id, commission_percent, flat_fee in (await self.db.execute(stmt)).all()
        }

    async def _resolve_bidder_profiles(
        self, event_id: UUID, bidder_numbers: set[int]
    ) -> dict[int, AuctioneerBidderProfile]:
        if not bidder_numbers:
            return {}

        stmt = (
            select(
                RegistrationGuest.bidder_number,
                RegistrationGuest.name,
                RegistrationGuest.table_number,
                RegistrationGuest.user_id,
                User.first_name,
                User.last_name,
                User.profile_picture_url,
            )
            .join(EventRegistration, EventRegistration.id == RegistrationGuest.registration_id)
            .outerjoin(User, User.id == RegistrationGuest.user_id)
            .where(
                EventRegistration.event_id == event_id,
                RegistrationGuest.bidder_number.in_(bidder_numbers),
            )
            .order_by(RegistrationGuest.is_primary.desc(), RegistrationGuest.created_at.asc())
        )
        rows = (await self.db.execute(stmt)).all()

        profiles: dict[int, AuctioneerBidderProfile] = {}
        for (
            bidder_number,
            guest_name,
            table_number,
            user_id,
            first_name,
            last_name,
            profile_picture_url,
        ) in rows:
            if bidder_number is None or bidder_number in profiles:
                continue

            bidder_name = (
                guest_name or " ".join(part for part in [first_name, last_name] if part).strip()
            )
            profiles[int(bidder_number)] = AuctioneerBidderProfile(
                bidder_number=int(bidder_number),
                bidder_name=bidder_name or "Unknown bidder",
                table_number=table_number,
                profile_picture_url=profile_picture_url,
                user_id=user_id,
            )

        return profiles

    async def _get_registered_bidder_count(self, event_id: UUID) -> int:
        stmt = (
            select(func.count(func.distinct(RegistrationGuest.bidder_number)))
            .join(EventRegistration, EventRegistration.id == RegistrationGuest.registration_id)
            .where(
                EventRegistration.event_id == event_id,
                RegistrationGuest.bidder_number.is_not(None),
            )
        )
        return int((await self.db.execute(stmt)).scalar_one() or 0)

    # ── Commission CRUD ──────────────────────────────────────────

    async def get_commissions(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> CommissionListResponse:
        # Subquery for primary image to avoid N+1 queries
        primary_img_subq = (
            select(
                AuctionItemMedia.auction_item_id,
                func.min(AuctionItemMedia.file_path).label("file_path"),
            )
            .where(AuctionItemMedia.display_order == 0)
            .group_by(AuctionItemMedia.auction_item_id)
            .subquery()
        )

        stmt = (
            select(
                AuctioneerItemCommission,
                AuctionItem.title,
                AuctionItem.bid_number,
                AuctionItem.auction_type,
                AuctionItem.status,
                AuctionItem.current_bid_amount,
                AuctionItem.quantity_available,
                AuctionItem.donor_value,
                AuctionItem.cost,
                AuctionItem.bid_count,
                primary_img_subq.c.file_path.label("primary_image_url"),
            )
            .join(
                AuctionItem,
                AuctioneerItemCommission.auction_item_id == AuctionItem.id,
            )
            .outerjoin(
                primary_img_subq,
                AuctionItem.id == primary_img_subq.c.auction_item_id,
            )
            .where(
                AuctioneerItemCommission.auctioneer_user_id == auctioneer_user_id,
                AuctionItem.event_id == event_id,
                AuctionItem.deleted_at.is_(None),
            )
            .order_by(AuctionItem.title)
        )
        rows = (await self.db.execute(stmt)).all()

        commissions: list[CommissionListItem] = []
        for row in rows:
            comm = row[0]
            primary_image = row[10]  # primary_image_url from subquery

            commissions.append(
                CommissionListItem(
                    id=comm.id,
                    auction_item_id=comm.auction_item_id,
                    auction_item_title=row[1],
                    auction_item_bid_number=row[2],
                    auction_type=row[3],
                    commission_percent=comm.commission_percent,
                    flat_fee=comm.flat_fee,
                    notes=comm.notes,
                    item_status=row[4],
                    current_bid_amount=row[5],
                    quantity_available=row[6],
                    donor_value=row[7],
                    cost=row[8],
                    bid_count=row[9],
                    primary_image_url=primary_image,
                    created_at=comm.created_at,
                    updated_at=comm.updated_at,
                )
            )

        return CommissionListResponse(commissions=commissions, total=len(commissions))

    async def upsert_commission(
        self,
        auctioneer_user_id: UUID,
        auction_item_id: UUID,
        commission_percent: Decimal,
        flat_fee: Decimal,
        notes: str | None,
    ) -> tuple[CommissionResponse, bool]:
        """Upsert a commission record. Returns (response, created)."""
        stmt = select(AuctioneerItemCommission).where(
            AuctioneerItemCommission.auctioneer_user_id == auctioneer_user_id,
            AuctioneerItemCommission.auction_item_id == auction_item_id,
        )
        existing = (await self.db.execute(stmt)).scalar_one_or_none()

        if existing:
            existing.commission_percent = commission_percent
            existing.flat_fee = flat_fee
            existing.notes = notes
            existing.updated_at = datetime.now(UTC)
            await self.db.flush()
            logger.info(
                "Commission updated: item %s by auctioneer %s (pct=%s, flat=%s)",
                auction_item_id,
                auctioneer_user_id,
                commission_percent,
                flat_fee,
            )
            return CommissionResponse.model_validate(existing), False

        new_comm = AuctioneerItemCommission(
            auctioneer_user_id=auctioneer_user_id,
            auction_item_id=auction_item_id,
            commission_percent=commission_percent,
            flat_fee=flat_fee,
            notes=notes,
        )
        self.db.add(new_comm)
        await self.db.flush()
        logger.info(
            "Commission created: item %s by auctioneer %s (pct=%s, flat=%s)",
            auction_item_id,
            auctioneer_user_id,
            commission_percent,
            flat_fee,
        )
        return CommissionResponse.model_validate(new_comm), True

    async def delete_commission(self, auctioneer_user_id: UUID, auction_item_id: UUID) -> bool:
        stmt = delete(AuctioneerItemCommission).where(
            AuctioneerItemCommission.auctioneer_user_id == auctioneer_user_id,
            AuctioneerItemCommission.auction_item_id == auction_item_id,
        )
        cursor: Any = await self.db.execute(stmt)
        deleted = bool(cursor.rowcount)
        if deleted:
            logger.info(
                "Commission deleted: item %s by auctioneer %s",
                auction_item_id,
                auctioneer_user_id,
            )
        return deleted

    # ── Event settings ───────────────────────────────────────────

    async def get_event_settings(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> EventSettingsResponse:
        stmt = select(AuctioneerEventSettings).where(
            AuctioneerEventSettings.auctioneer_user_id == auctioneer_user_id,
            AuctioneerEventSettings.event_id == event_id,
        )
        settings = (await self.db.execute(stmt)).scalar_one_or_none()

        if settings:
            return EventSettingsResponse.model_validate(settings)

        # Return defaults
        return EventSettingsResponse(
            auctioneer_user_id=auctioneer_user_id,
            event_id=event_id,
            live_auction_percent=Decimal("0"),
            paddle_raise_percent=Decimal("0"),
            silent_auction_percent=Decimal("0"),
            paddle_raise_levels=DEFAULT_PADDLE_RAISE_LEVELS.copy(),
            paddle_raise_total_goal=None,
            paddle_raise_level_goals={},
            paddle_raise_level_notes={},
            created_at=datetime.now(UTC),
            updated_at=datetime.now(UTC),
        )

    async def upsert_event_settings(
        self,
        event_id: UUID,
        auctioneer_user_id: UUID,
        live_auction_percent: Decimal,
        paddle_raise_percent: Decimal,
        silent_auction_percent: Decimal,
        paddle_raise_levels: list[int],
        paddle_raise_total_goal: int | None,
        paddle_raise_level_goals: dict[str, int],
        paddle_raise_level_notes: dict[str, str],
    ) -> EventSettingsResponse:
        normalized_total_goal = (
            Decimal(paddle_raise_total_goal) if paddle_raise_total_goal is not None else None
        )
        stmt = select(AuctioneerEventSettings).where(
            AuctioneerEventSettings.auctioneer_user_id == auctioneer_user_id,
            AuctioneerEventSettings.event_id == event_id,
        )
        existing = (await self.db.execute(stmt)).scalar_one_or_none()

        if existing:
            existing.live_auction_percent = live_auction_percent
            existing.paddle_raise_percent = paddle_raise_percent
            existing.silent_auction_percent = silent_auction_percent
            existing.paddle_raise_levels = paddle_raise_levels
            existing.paddle_raise_total_goal = normalized_total_goal
            existing.paddle_raise_level_goals = paddle_raise_level_goals
            existing.paddle_raise_level_notes = paddle_raise_level_notes
            existing.updated_at = datetime.now(UTC)
            await self.db.flush()
            logger.info(
                "Event settings updated: event %s by auctioneer %s",
                event_id,
                auctioneer_user_id,
            )
            return EventSettingsResponse.model_validate(existing)

        new_settings = AuctioneerEventSettings(
            auctioneer_user_id=auctioneer_user_id,
            event_id=event_id,
            live_auction_percent=live_auction_percent,
            paddle_raise_percent=paddle_raise_percent,
            silent_auction_percent=silent_auction_percent,
            paddle_raise_levels=paddle_raise_levels,
            paddle_raise_total_goal=normalized_total_goal,
            paddle_raise_level_goals=paddle_raise_level_goals,
            paddle_raise_level_notes=paddle_raise_level_notes,
        )
        self.db.add(new_settings)
        await self.db.flush()
        logger.info(
            "Event settings created: event %s by auctioneer %s",
            event_id,
            auctioneer_user_id,
        )
        return EventSettingsResponse.model_validate(new_settings)

    # ── Dashboard / earnings ─────────────────────────────────────

    async def _get_event_revenue(self, event_id: UUID) -> dict[str, Decimal]:
        """Get event revenue totals by category (same pattern as EventDashboardService)."""
        # Silent auction high bids
        silent_high = (
            select(
                AuctionBid.auction_item_id,
                func.max(AuctionBid.bid_amount).label("max_bid"),
            )
            .join(AuctionItem, AuctionBid.auction_item_id == AuctionItem.id)
            .where(
                AuctionBid.event_id == event_id,
                AuctionBid.bid_status.in_(
                    [BidStatus.ACTIVE.value, BidStatus.OUTBID.value, BidStatus.WINNING.value]
                ),
                AuctionItem.auction_type == "silent",
                AuctionItem.deleted_at.is_(None),
            )
            .group_by(AuctionBid.auction_item_id)
            .subquery()
        )
        silent_stmt = select(func.coalesce(func.sum(silent_high.c.max_bid), 0))
        silent_total = Decimal((await self.db.execute(silent_stmt)).scalar_one() or 0)

        # Live auction - QuickEntryBid (capped at buy-it-now price if sold via buy-now)
        buy_now_min_sq = (
            select(
                QuickEntryBuyNowBid.item_id,
                func.min(QuickEntryBuyNowBid.amount).label("min_price"),
            )
            .where(QuickEntryBuyNowBid.event_id == event_id)
            .group_by(QuickEntryBuyNowBid.item_id)
            .subquery("bn_min_ev")
        )
        live_high = (
            select(
                func.least(
                    func.max(QuickEntryBid.amount),
                    func.coalesce(buy_now_min_sq.c.min_price, func.max(QuickEntryBid.amount)),
                ).label("effective_bid"),
            )
            .outerjoin(buy_now_min_sq, QuickEntryBid.item_id == buy_now_min_sq.c.item_id)
            .where(
                QuickEntryBid.event_id == event_id,
                QuickEntryBid.status.in_([QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]),
            )
            .group_by(QuickEntryBid.item_id, buy_now_min_sq.c.min_price)
            .subquery("qe_live_high_ev")
        )
        live_stmt = select(func.coalesce(func.sum(live_high.c.effective_bid), 0))
        live_total = Decimal((await self.db.execute(live_stmt)).scalar_one() or 0)

        # Paddle raise
        paddle_stmt = select(func.coalesce(func.sum(PaddleRaiseContribution.amount), 0)).where(
            PaddleRaiseContribution.event_id == event_id
        )
        qe_paddle_stmt = select(func.coalesce(func.sum(QuickEntryDonation.amount), 0)).where(
            QuickEntryDonation.event_id == event_id
        )
        paddle_total = Decimal((await self.db.execute(paddle_stmt)).scalar_one() or 0)
        qe_paddle_total = Decimal((await self.db.execute(qe_paddle_stmt)).scalar_one() or 0)

        return {
            "live_auction": live_total,
            "paddle_raise": paddle_total + qe_paddle_total,
            "silent_auction": silent_total,
        }

    async def _get_per_item_earnings(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> tuple[Decimal, int, set[UUID]]:
        """Calculate per-item earnings. Returns (total, count, commissioned_item_ids)."""
        comm_stmt = (
            select(
                AuctioneerItemCommission.auction_item_id,
                AuctioneerItemCommission.commission_percent,
                AuctioneerItemCommission.flat_fee,
                AuctionItem.donor_value,
            )
            .join(
                AuctionItem,
                AuctioneerItemCommission.auction_item_id == AuctionItem.id,
            )
            .where(
                AuctioneerItemCommission.auctioneer_user_id == auctioneer_user_id,
                AuctionItem.event_id == event_id,
                AuctionItem.deleted_at.is_(None),
                AuctionItem.status.notin_(["withdrawn", "draft"]),
            )
        )
        comm_rows = (await self.db.execute(comm_stmt)).all()

        if not comm_rows:
            return Decimal("0"), 0, set()

        commissioned_ids = {row[0] for row in comm_rows}

        total = Decimal("0")
        count = 0
        for _item_id, pct, flat, donor_value in comm_rows:
            sale_amount = Decimal(donor_value or 0)
            earning = sale_amount * pct / Decimal("100") + Decimal(flat or 0)
            if earning > 0:
                total += earning
                count += 1

        return total, count, commissioned_ids

    async def get_dashboard(self, event_id: UUID, auctioneer_user_id: UUID) -> DashboardResponse:
        now = datetime.now(UTC)

        # Get event for timer data
        event = (
            await self.db.execute(select(Event).where(Event.id == event_id))
        ).scalar_one_or_none()

        # Per-item earnings
        per_item_total, per_item_count, commissioned_ids = await self._get_per_item_earnings(
            event_id, auctioneer_user_id
        )

        # Event revenue totals
        revenue = await self._get_event_revenue(event_id)

        # Event settings for category percentages
        settings = await self.get_event_settings(event_id, auctioneer_user_id)

        # Category earnings — exclude commissioned items to avoid double-counting.
        # Category percentages apply only to revenue from items without per-item commissions.
        commissioned_live_revenue = Decimal("0")
        commissioned_silent_revenue = Decimal("0")

        if commissioned_ids:
            for c_id in commissioned_ids:
                bid_amt = Decimal("0")
                # Check auction type for this item
                item_type_result = await self.db.execute(
                    select(AuctionItem.auction_type).where(AuctionItem.id == c_id)
                )
                item_type = item_type_result.scalar_one_or_none()

                # Get highest bid from AuctionBid
                ab_stmt = select(func.max(AuctionBid.bid_amount)).where(
                    AuctionBid.auction_item_id == c_id,
                    AuctionBid.bid_status.in_(
                        [BidStatus.ACTIVE.value, BidStatus.OUTBID.value, BidStatus.WINNING.value]
                    ),
                )
                ab_max = Decimal((await self.db.execute(ab_stmt)).scalar_one() or 0)

                # Get highest bid from QuickEntryBid
                qe_stmt = select(func.max(QuickEntryBid.amount)).where(
                    QuickEntryBid.item_id == c_id,
                    QuickEntryBid.status.in_(
                        [QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]
                    ),
                )
                qe_max = Decimal((await self.db.execute(qe_stmt)).scalar_one() or 0)

                bid_amt = max(ab_max, qe_max)
                if item_type == "live":
                    commissioned_live_revenue += bid_amt
                elif item_type == "silent":
                    commissioned_silent_revenue += bid_amt

        live_pool = max(Decimal("0"), revenue["live_auction"] - commissioned_live_revenue)
        silent_pool = max(Decimal("0"), revenue["silent_auction"] - commissioned_silent_revenue)

        live_cat = live_pool * settings.live_auction_percent / Decimal("100")
        paddle_cat = revenue["paddle_raise"] * settings.paddle_raise_percent / Decimal("100")
        silent_cat = silent_pool * settings.silent_auction_percent / Decimal("100")

        total_earnings = per_item_total + live_cat + paddle_cat + silent_cat

        event_total_raised = (
            revenue["live_auction"] + revenue["paddle_raise"] + revenue["silent_auction"]
        )

        # Timer calculations
        live_start = getattr(event, "live_auction_start_datetime", None) if event else None
        auction_close = getattr(event, "auction_close_datetime", None) if event else None

        live_status: AuctionStatus = "not_scheduled"
        if live_start:
            if now < live_start:
                live_status = "not_started"
            else:
                # Check if any live items are still bidding_open
                active_live = (
                    await self.db.execute(
                        select(func.count())
                        .select_from(AuctionItem)
                        .where(
                            AuctionItem.event_id == event_id,
                            AuctionItem.auction_type == "live",
                            AuctionItem.bidding_open.is_(True),
                            AuctionItem.deleted_at.is_(None),
                        )
                    )
                ).scalar_one()
                live_status = "in_progress" if active_live > 0 else "ended"

        silent_status: SilentAuctionStatus = "not_scheduled"
        if auction_close:
            event_start = event.event_datetime if event else None
            if now >= auction_close:
                silent_status = "closed"
            elif event_start and now < event_start:
                silent_status = "not_started"
            else:
                silent_status = "open"

        return DashboardResponse(
            earnings=EarningsSummary(
                per_item_total=per_item_total,
                per_item_count=per_item_count,
                live_auction_category_earning=live_cat,
                paddle_raise_category_earning=paddle_cat,
                silent_auction_category_earning=silent_cat,
                total_earnings=total_earnings,
            ),
            event_totals=EventTotals(
                live_auction_raised=revenue["live_auction"],
                paddle_raise_raised=revenue["paddle_raise"],
                silent_auction_raised=revenue["silent_auction"],
                event_total_raised=event_total_raised,
            ),
            timers=TimerData(
                live_auction_start_datetime=live_start,
                auction_close_datetime=auction_close,
                live_auction_status=live_status,
                silent_auction_status=silent_status,
            ),
            last_refreshed_at=now,
            revenue_generators=await self._revenue_generators_summary(event_id),
        )

    async def _revenue_generators_summary(
        self, event_id: UUID
    ) -> RevenueGeneratorDashboardSummary | None:
        try:
            from app.services.revenue_generator_service import RevenueGeneratorService

            result = await RevenueGeneratorService.get_dashboard_summary(self.db, event_id)
            return result
        except Exception:
            logger.exception("Failed to load revenue generator summary for event %s", event_id)
            return None

    # ── Live auction tab ─────────────────────────────────────────

    async def get_live_auction(self, event_id: UUID) -> LiveAuctionResponse:
        """Get current live auction item, high bidder, and bid history."""
        # Find current live item (bidding_open + live type)
        item_stmt = (
            select(AuctionItem)
            .where(
                AuctionItem.event_id == event_id,
                AuctionItem.auction_type == "live",
                AuctionItem.bidding_open.is_(True),
                AuctionItem.deleted_at.is_(None),
            )
            .order_by(AuctionItem.updated_at.desc())
            .limit(1)
        )
        current_item = (await self.db.execute(item_stmt)).scalar_one_or_none()

        if not current_item:
            # Check if auction has started at all
            any_live = (
                await self.db.execute(
                    select(func.count())
                    .select_from(AuctionItem)
                    .where(
                        AuctionItem.event_id == event_id,
                        AuctionItem.auction_type == "live",
                        AuctionItem.deleted_at.is_(None),
                    )
                )
            ).scalar_one()

            return LiveAuctionResponse(
                current_item=None,
                high_bidder=None,
                bid_history=[],
                auction_status="ended" if any_live > 0 else "not_started",
            )

        # Get primary image
        img_stmt = (
            select(AuctionItemMedia.file_path)
            .where(
                AuctionItemMedia.auction_item_id == current_item.id,
                AuctionItemMedia.display_order == 0,
            )
            .limit(1)
        )
        primary_image = (await self.db.execute(img_stmt)).scalar_one_or_none()

        live_item = LiveAuctionItem(
            id=current_item.id,
            bid_number=current_item.bid_number,
            title=current_item.title,
            description=current_item.description,
            starting_bid=current_item.starting_bid,
            current_bid_amount=current_item.current_bid_amount,
            bid_count=current_item.bid_count or 0,
            primary_image_url=primary_image,
            donor_value=current_item.donor_value,
            cost=current_item.cost,
        )

        # Get high bidder from quick entry bids (live auction uses QuickEntryBid)
        high_bid_stmt = (
            select(QuickEntryBid)
            .where(
                QuickEntryBid.item_id == current_item.id,
                QuickEntryBid.status.in_([QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]),
            )
            .order_by(QuickEntryBid.amount.desc())
            .limit(1)
        )
        high_bid = (await self.db.execute(high_bid_stmt)).scalar_one_or_none()

        high_bidder = None
        if high_bid:
            # Get user + registration guest details
            user_stmt = select(User).where(User.id == high_bid.donor_user_id)
            bidder_user = (await self.db.execute(user_stmt)).scalar_one_or_none()

            table_number = None
            if bidder_user:
                guest_stmt = (
                    select(RegistrationGuest.table_number)
                    .join(
                        EventRegistration,
                        RegistrationGuest.registration_id == EventRegistration.id,
                    )
                    .where(
                        RegistrationGuest.user_id == bidder_user.id,
                        EventRegistration.event_id == event_id,
                    )
                    .limit(1)
                )
                guest = (await self.db.execute(guest_stmt)).scalar_one_or_none()
                table_number = guest

            high_bidder = HighBidder(
                bidder_number=high_bid.bidder_number,
                first_name=bidder_user.first_name if bidder_user else "Unknown",
                last_name=bidder_user.last_name if bidder_user else "",
                table_number=table_number,
                profile_picture_url=(bidder_user.profile_picture_url if bidder_user else None),
            )

        # Bid history from QuickEntryBid
        history_stmt = (
            select(QuickEntryBid, User.first_name, User.last_name)
            .outerjoin(User, QuickEntryBid.donor_user_id == User.id)
            .where(
                QuickEntryBid.item_id == current_item.id,
                QuickEntryBid.status.in_([QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]),
            )
            .order_by(QuickEntryBid.amount.desc())
        )
        history_rows = (await self.db.execute(history_stmt)).all()

        bid_history = [
            BidHistoryEntry(
                bidder_number=row[0].bidder_number,
                bidder_name=f"{row[1] or 'Unknown'} {row[2] or ''}".strip(),
                bid_amount=Decimal(row[0].amount),
                placed_at=row[0].accepted_at or row[0].created_at,
            )
            for row in history_rows
        ]

        return LiveAuctionResponse(
            current_item=live_item,
            high_bidder=high_bidder,
            bid_history=bid_history,
            auction_status="in_progress",
        )

    async def get_live_auction_gallery(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> AuctioneerItemGalleryResponse:
        item_stmt = (
            select(AuctionItem)
            .where(
                AuctionItem.event_id == event_id,
                AuctionItem.auction_type == "live",
                AuctionItem.deleted_at.is_(None),
            )
            .order_by(AuctionItem.status.asc(), AuctionItem.bid_number.asc())
        )
        items = list((await self.db.execute(item_stmt)).scalars().all())
        item_ids = [item.id for item in items]

        bid_stats = {
            item_id: (Decimal(max_bid or 0), int(bid_count or 0), int(bidder_count or 0))
            for item_id, max_bid, bid_count, bidder_count in (
                await self.db.execute(
                    select(
                        QuickEntryBid.item_id,
                        func.max(QuickEntryBid.amount),
                        func.count(QuickEntryBid.id),
                        func.count(func.distinct(QuickEntryBid.bidder_number)),
                    )
                    .where(
                        QuickEntryBid.event_id == event_id,
                        QuickEntryBid.status.in_(
                            [QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]
                        ),
                    )
                    .group_by(QuickEntryBid.item_id)
                )
            ).all()
        }
        buy_now_min_prices: dict[UUID, Decimal] = {
            row.item_id: Decimal(row.min_price)
            for row in (
                await self.db.execute(
                    select(
                        QuickEntryBuyNowBid.item_id,
                        func.min(QuickEntryBuyNowBid.amount).label("min_price"),
                    )
                    .where(QuickEntryBuyNowBid.event_id == event_id)
                    .group_by(QuickEntryBuyNowBid.item_id)
                )
            ).all()
        }
        commission_map = await self._get_commission_map(event_id, auctioneer_user_id)
        image_map = await self._get_primary_image_map(item_ids)

        summaries: list[AuctioneerItemSummary] = []
        total_raised = Decimal("0")
        total_bids = 0

        for item in items:
            current_bid_amount, bid_count, bidder_count = bid_stats.get(
                item.id, (Decimal("0"), 0, 0)
            )
            if item.id in buy_now_min_prices:
                current_bid_amount = min(current_bid_amount, buy_now_min_prices[item.id])
            total_raised += current_bid_amount
            total_bids += bid_count
            commission_data = commission_map.get(item.id)
            commission_percent = commission_data[0] if commission_data is not None else None
            flat_fee = commission_data[1] if commission_data is not None else None
            summaries.append(
                AuctioneerItemSummary(
                    id=item.id,
                    bid_number=item.bid_number,
                    title=item.title,
                    auction_type=item.auction_type,
                    status=item.status,
                    description=item.description,
                    current_bid_amount=current_bid_amount if current_bid_amount > 0 else None,
                    bid_count=bid_count,
                    bidder_count=bidder_count,
                    primary_image_url=image_map.get(item.id),
                    donor_value=item.donor_value,
                    cost=item.cost,
                    commission_percent=commission_percent,
                    flat_fee=flat_fee,
                    has_commission=commission_percent is not None,
                    has_bounty=bool(flat_fee and flat_fee > 0),
                    slide_presentation_html=item.slide_presentation_html,
                    slide_presentation_layout=item.slide_presentation_layout,
                )
            )

        return AuctioneerItemGalleryResponse(
            items=summaries,
            total_items=len(summaries),
            total_raised=total_raised,
            total_bids=total_bids,
        )

    async def get_silent_auction_gallery(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> AuctioneerItemGalleryResponse:
        item_stmt = (
            select(AuctionItem)
            .where(
                AuctionItem.event_id == event_id,
                AuctionItem.auction_type == "silent",
                AuctionItem.deleted_at.is_(None),
            )
            .order_by(AuctionItem.status.asc(), AuctionItem.bid_number.asc())
        )
        items = list((await self.db.execute(item_stmt)).scalars().all())
        item_ids = [item.id for item in items]

        bid_stats = {
            item_id: (Decimal(max_bid or 0), int(bid_count or 0), int(bidder_count or 0))
            for item_id, max_bid, bid_count, bidder_count in (
                await self.db.execute(
                    select(
                        AuctionBid.auction_item_id,
                        func.max(AuctionBid.bid_amount),
                        func.count(AuctionBid.id),
                        func.count(func.distinct(AuctionBid.bidder_number)),
                    )
                    .where(
                        AuctionBid.event_id == event_id,
                        AuctionBid.bid_status.notin_(
                            [BidStatus.CANCELLED.value, BidStatus.WITHDRAWN.value]
                        ),
                    )
                    .group_by(AuctionBid.auction_item_id)
                )
            ).all()
        }
        commission_map = await self._get_commission_map(event_id, auctioneer_user_id)
        image_map = await self._get_primary_image_map(item_ids)

        summaries: list[AuctioneerItemSummary] = []
        total_raised = Decimal("0")
        total_bids = 0

        for item in items:
            current_bid_amount, bid_count, bidder_count = bid_stats.get(
                item.id, (Decimal("0"), 0, 0)
            )
            total_raised += current_bid_amount
            total_bids += bid_count
            commission_data = commission_map.get(item.id)
            commission_percent = commission_data[0] if commission_data is not None else None
            flat_fee = commission_data[1] if commission_data is not None else None
            summaries.append(
                AuctioneerItemSummary(
                    id=item.id,
                    bid_number=item.bid_number,
                    title=item.title,
                    auction_type=item.auction_type,
                    status=item.status,
                    description=item.description,
                    current_bid_amount=current_bid_amount if current_bid_amount > 0 else None,
                    bid_count=bid_count,
                    bidder_count=bidder_count,
                    primary_image_url=image_map.get(item.id),
                    donor_value=item.donor_value,
                    cost=item.cost,
                    commission_percent=commission_percent,
                    flat_fee=flat_fee,
                    has_commission=commission_percent is not None,
                    has_bounty=bool(flat_fee and flat_fee > 0),
                    slide_presentation_html=item.slide_presentation_html,
                    slide_presentation_layout=item.slide_presentation_layout,
                )
            )

        return AuctioneerItemGalleryResponse(
            items=summaries,
            total_items=len(summaries),
            total_raised=total_raised,
            total_bids=total_bids,
        )

    async def get_item_detail(
        self, event_id: UUID, item_id: UUID, auctioneer_user_id: UUID
    ) -> AuctioneerItemDetailResponse:
        item = (
            await self.db.execute(
                select(AuctionItem).where(
                    AuctionItem.id == item_id,
                    AuctionItem.event_id == event_id,
                    AuctionItem.deleted_at.is_(None),
                )
            )
        ).scalar_one_or_none()
        if item is None:
            raise ValueError("Auction item not found")

        commission_map = await self._get_commission_map(event_id, auctioneer_user_id)
        image_map = await self._get_primary_image_map([item.id])
        commission_data = commission_map.get(item.id)
        commission_percent = commission_data[0] if commission_data is not None else None
        flat_fee = commission_data[1] if commission_data is not None else None

        if item.auction_type == "live":
            live_bid_rows = (
                await self.db.execute(
                    select(
                        QuickEntryBid.id,
                        QuickEntryBid.amount,
                        QuickEntryBid.accepted_at,
                        QuickEntryBid.status,
                        QuickEntryBid.bidder_number,
                    )
                    .where(
                        QuickEntryBid.event_id == event_id,
                        QuickEntryBid.item_id == item_id,
                        QuickEntryBid.status.in_(
                            [QuickEntryBidStatus.ACTIVE, QuickEntryBidStatus.WINNING]
                        ),
                    )
                    .order_by(QuickEntryBid.accepted_at.desc(), QuickEntryBid.amount.desc())
                )
            ).all()
            bidder_profiles = await self._resolve_bidder_profiles(
                event_id,
                {int(row.bidder_number) for row in live_bid_rows if row.bidder_number is not None},
            )

            bids = [
                AuctioneerBidActivityEntry(
                    id=row.id,
                    bid_amount=Decimal(row.amount),
                    placed_at=row.accepted_at,
                    bid_status=row.status.value,
                    bid_source="live",
                    bidder=bidder_profiles.get(
                        int(row.bidder_number),
                        AuctioneerBidderProfile(
                            bidder_number=row.bidder_number,
                            bidder_name=f"Bidder #{row.bidder_number}",
                            table_number=None,
                            profile_picture_url=None,
                        ),
                    ),
                )
                for row in live_bid_rows
            ]
        else:
            silent_bid_rows = (
                await self.db.execute(
                    select(
                        AuctionBid.id,
                        AuctionBid.bid_amount,
                        AuctionBid.placed_at,
                        AuctionBid.bid_status,
                        AuctionBid.bidder_number,
                        AuctionBid.user_id,
                        User.first_name,
                        User.last_name,
                        User.profile_picture_url,
                    )
                    .outerjoin(User, User.id == AuctionBid.user_id)
                    .where(
                        AuctionBid.event_id == event_id,
                        AuctionBid.auction_item_id == item_id,
                        AuctionBid.bid_status.notin_(
                            [BidStatus.CANCELLED.value, BidStatus.WITHDRAWN.value]
                        ),
                    )
                    .order_by(AuctionBid.placed_at.desc(), AuctionBid.bid_amount.desc())
                )
            ).all()
            bidder_profiles = await self._resolve_bidder_profiles(
                event_id,
                {
                    int(row.bidder_number)
                    for row in silent_bid_rows
                    if row.bidder_number is not None
                },
            )

            bids = []
            for row in silent_bid_rows:
                fallback_name = (
                    " ".join(part for part in [row.first_name, row.last_name] if part).strip()
                    or f"Bidder #{row.bidder_number}"
                )
                bidder = bidder_profiles.get(
                    int(row.bidder_number),
                    AuctioneerBidderProfile(
                        bidder_number=row.bidder_number,
                        bidder_name=fallback_name,
                        table_number=None,
                        profile_picture_url=row.profile_picture_url,
                        user_id=row.user_id,
                    ),
                )
                bids.append(
                    AuctioneerBidActivityEntry(
                        id=row.id,
                        bid_amount=Decimal(row.bid_amount),
                        placed_at=row.placed_at,
                        bid_status=row.bid_status,
                        bid_source="silent",
                        bidder=bidder,
                    )
                )

        grouped: dict[int, list[AuctioneerBidActivityEntry]] = defaultdict(list)
        for bid in bids:
            if bid.bidder.bidder_number is not None:
                grouped[bid.bidder.bidder_number].append(bid)

        bidder_summaries = [
            AuctioneerBidderSummary(
                bidder=entries[0].bidder,
                total_bid_amount=sum((entry.bid_amount for entry in entries), Decimal("0")),
                highest_bid_amount=max(entry.bid_amount for entry in entries),
                bid_count=len(entries),
                latest_bid_at=max(entry.placed_at for entry in entries),
            )
            for entries in grouped.values()
        ]
        bidder_summaries.sort(
            key=lambda summary: (-summary.highest_bid_amount, -summary.total_bid_amount)
        )

        current_high_bid = bidder_summaries[0].highest_bid_amount if bidder_summaries else None
        high_bidder = bidder_summaries[0].bidder if bidder_summaries else None

        return AuctioneerItemDetailResponse(
            item=AuctioneerItemSummary(
                id=item.id,
                bid_number=item.bid_number,
                title=item.title,
                auction_type=item.auction_type,
                status=item.status,
                description=item.description,
                current_bid_amount=current_high_bid,
                bid_count=len(bids),
                bidder_count=len(grouped),
                primary_image_url=image_map.get(item.id),
                donor_value=item.donor_value,
                cost=item.cost,
                commission_percent=commission_percent,
                flat_fee=flat_fee,
                has_commission=commission_percent is not None,
                has_bounty=bool(flat_fee and flat_fee > 0),
                slide_presentation_html=item.slide_presentation_html,
                slide_presentation_layout=item.slide_presentation_layout,
            ),
            current_high_bid=current_high_bid,
            high_bidder=high_bidder,
            bids=bids,
            bidder_summaries=bidder_summaries,
        )

    async def get_paddle_raise(
        self, event_id: UUID, auctioneer_user_id: UUID
    ) -> AuctioneerPaddleRaiseResponse:
        settings = await self.get_event_settings(event_id, auctioneer_user_id)
        donations = list(
            (
                await self.db.execute(
                    select(QuickEntryDonation)
                    .where(QuickEntryDonation.event_id == event_id)
                    .options(
                        selectinload(QuickEntryDonation.label_links).selectinload(
                            QuickEntryDonationLabelLink.label
                        )
                    )
                    .order_by(QuickEntryDonation.entered_at.desc())
                )
            )
            .scalars()
            .all()
        )

        bidder_profiles = await self._resolve_bidder_profiles(
            event_id,
            {int(donation.bidder_number) for donation in donations},
        )
        registered_bidder_count = await self._get_registered_bidder_count(event_id)

        donation_entries: list[AuctioneerBidActivityEntry] = []
        level_totals: dict[tuple[int, bool], Decimal] = defaultdict(lambda: Decimal("0"))
        level_bidder_sets: dict[tuple[int, bool], set[int]] = defaultdict(set)
        bidder_total_map: dict[int, Decimal] = defaultdict(lambda: Decimal("0"))
        bidder_count_map: dict[int, int] = defaultdict(int)
        bidder_labels_map: dict[int, set[str]] = defaultdict(set)
        last_hero_total_map: dict[int, Decimal] = defaultdict(lambda: Decimal("0"))
        last_hero_count_map: dict[int, int] = defaultdict(int)
        last_hero_labels_map: dict[int, set[str]] = defaultdict(set)
        last_hero_total = Decimal("0")

        for donation in donations:
            bidder = bidder_profiles.get(
                donation.bidder_number,
                AuctioneerBidderProfile(
                    bidder_number=donation.bidder_number,
                    bidder_name=f"Bidder #{donation.bidder_number}",
                    table_number=None,
                    profile_picture_url=None,
                ),
            )
            label_names = [
                self._normalize_donation_label(link.label.name)
                if link.label is not None
                else self._normalize_donation_label(link.custom_label_text or "")
                for link in donation.label_links
                if link.label is not None or link.custom_label_text
            ]
            is_last_hero = any(label.lower() == "last hero" for label in label_names)

            donation_entries.append(
                AuctioneerBidActivityEntry(
                    id=donation.id,
                    bid_amount=Decimal(donation.amount),
                    placed_at=donation.entered_at,
                    bid_status="recorded",
                    bid_source="paddle_raise",
                    label_names=label_names,
                    is_monthly=donation.is_monthly,
                    bidder=bidder,
                )
            )

            level_key = (int(donation.amount), bool(donation.is_monthly))
            level_totals[level_key] += Decimal(donation.amount)
            level_bidder_sets[level_key].add(donation.bidder_number)

            bidder_total_map[donation.bidder_number] += Decimal(donation.amount)
            bidder_count_map[donation.bidder_number] += 1
            bidder_labels_map[donation.bidder_number].update(label_names)

            if is_last_hero:
                last_hero_total += Decimal(donation.amount)
                last_hero_total_map[donation.bidder_number] += Decimal(donation.amount)
                last_hero_count_map[donation.bidder_number] += 1
                last_hero_labels_map[donation.bidder_number].update(label_names)

        configured_keys = [(amount, False) for amount in settings.paddle_raise_levels]
        dynamic_keys = [key for key in level_totals.keys() if key not in configured_keys]
        level_keys = configured_keys + sorted(dynamic_keys, reverse=True)

        level_summaries = []
        for amount, is_monthly in level_keys:
            bidder_count = len(level_bidder_sets.get((amount, is_monthly), set()))
            participation_percent = (
                round((bidder_count / registered_bidder_count) * 100, 2)
                if registered_bidder_count > 0
                else 0.0
            )
            donations_count = sum(
                1
                for donation in donations
                if donation.amount == amount and bool(donation.is_monthly) == is_monthly
            )
            goal_amount = settings.paddle_raise_level_goals.get(str(amount))
            goal_progress_percent = (
                round(
                    (float(level_totals.get((amount, is_monthly), Decimal("0"))) / goal_amount)
                    * 100,
                    2,
                )
                if goal_amount
                else None
            )
            level_summaries.append(
                AuctioneerPaddleRaiseLevelSummary(
                    amount=amount,
                    bidder_count=bidder_count,
                    total_amount=level_totals.get((amount, is_monthly), Decimal("0")),
                    participation_percent=participation_percent,
                    donations_count=donations_count,
                    goal_amount=Decimal(goal_amount) if goal_amount else None,
                    goal_progress_percent=goal_progress_percent,
                    is_monthly=is_monthly,
                )
            )

        total_pledged = sum((Decimal(donation.amount) for donation in donations), Decimal("0"))
        unique_donor_count = len({donation.bidder_number for donation in donations})
        participation_percent = (
            round((unique_donor_count / registered_bidder_count) * 100, 2)
            if registered_bidder_count > 0
            else 0.0
        )
        total_goal = settings.paddle_raise_total_goal
        total_goal_progress_percent = (
            round((float(total_pledged) / float(total_goal)) * 100, 2) if total_goal else None
        )

        bidder_totals = [
            AuctioneerPaddleRaiseBidderSummary(
                bidder=bidder_profiles.get(
                    bidder_number,
                    AuctioneerBidderProfile(
                        bidder_number=bidder_number,
                        bidder_name=f"Bidder #{bidder_number}",
                        table_number=None,
                        profile_picture_url=None,
                    ),
                ),
                total_amount=total_amount,
                donation_count=bidder_count_map[bidder_number],
                label_names=sorted(bidder_labels_map[bidder_number]),
                is_last_hero=bidder_number in last_hero_total_map,
            )
            for bidder_number, total_amount in bidder_total_map.items()
        ]
        bidder_totals.sort(key=lambda row: (-row.total_amount, -row.donation_count))

        last_hero_bidder_totals = [
            AuctioneerPaddleRaiseBidderSummary(
                bidder=bidder_profiles.get(
                    bidder_number,
                    AuctioneerBidderProfile(
                        bidder_number=bidder_number,
                        bidder_name=f"Bidder #{bidder_number}",
                        table_number=None,
                        profile_picture_url=None,
                    ),
                ),
                total_amount=total_amount,
                donation_count=last_hero_count_map[bidder_number],
                label_names=sorted(last_hero_labels_map[bidder_number]),
                is_last_hero=True,
            )
            for bidder_number, total_amount in last_hero_total_map.items()
        ]
        last_hero_bidder_totals.sort(key=lambda row: (-row.total_amount, -row.donation_count))

        return AuctioneerPaddleRaiseResponse(
            configured_levels=settings.paddle_raise_levels,
            total_pledged=total_pledged,
            total_goal=total_goal,
            total_goal_progress_percent=total_goal_progress_percent,
            donation_count=len(donations),
            unique_donor_count=unique_donor_count,
            participation_percent=participation_percent,
            last_hero_total=last_hero_total,
            level_summaries=level_summaries,
            donations=donation_entries,
            bidder_totals=bidder_totals,
            last_hero_bidder_totals=last_hero_bidder_totals,
        )

    async def build_slide_deck(
        self, items: list[AuctioneerItemSummary], event_name: str, auction_type_label: str
    ) -> bytes:
        from io import BytesIO

        from PIL import Image
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
            for item in items:
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])

                # Extract the slide text content
                body_text = self._strip_slide_html(item.slide_presentation_html) or (
                    item.description or "No description provided."
                )

                layout_key = item.slide_presentation_layout

                # Layout configurations matching frontend SlidePreview component
                # Each layout: (image_area, text_area)
                layout_configs = {
                    # ON_IMAGE: Image fills most of slide, text overlay at bottom
                    "on_image": {
                        "image": (Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5)),
                        "title": (Inches(0.8), Inches(5.8), Inches(11.733), Inches(0.7)),
                        "text": (Inches(0.8), Inches(6.5), Inches(11.733), Inches(1.3)),
                    },
                    # LEFT_OF_IMAGE: Text on left, image on right (45%/55% split)
                    "left_of_image": {
                        "title": (Inches(0.5), Inches(0.5), Inches(5.5), Inches(0.6)),
                        "text": (Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.7)),
                        "image": (Inches(6.5), Inches(0.5), Inches(6.333), Inches(6.5)),
                    },
                    # RIGHT_OF_IMAGE: Image on left, text on right (55%/45% split)
                    "right_of_image": {
                        "image": (Inches(0.5), Inches(0.5), Inches(6.5), Inches(6.5)),
                        "title": (Inches(7.5), Inches(0.5), Inches(5.333), Inches(0.6)),
                        "text": (Inches(7.5), Inches(1.3), Inches(5.333), Inches(5.7)),
                    },
                    # BELOW_IMAGE: Image on top, text below
                    "below_image": {
                        "image": (Inches(0.5), Inches(0.5), Inches(12.333), Inches(4.5)),
                        "title": (Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.5)),
                        "text": (Inches(0.5), Inches(6.0), Inches(12.333), Inches(1.0)),
                    },
                }

                config = layout_configs[layout_key]

                # Add image if available
                if item.primary_image_url:
                    try:
                        logger.info(
                            "Fetching slide image for item %s (%s) from URL: %s",
                            item.id,
                            item.title,
                            item.primary_image_url[:100],
                        )
                        response = await client.get(item.primary_image_url)
                        response.raise_for_status()
                        image_bytes = response.content

                        if not image_bytes:
                            logger.warning(
                                "Empty image response for auction item %s (%s)",
                                item.id,
                                item.title,
                            )
                            raise ValueError("Empty image content")

                        # Get image dimensions to calculate aspect-ratio-preserving size
                        pil_image = Image.open(BytesIO(image_bytes))
                        img_width, img_height = pil_image.size
                        aspect_ratio = img_width / img_height

                        # Get max bounds for this layout
                        max_left, max_top, max_width, max_height = config["image"]

                        # Calculate dimensions that fit within bounds while maintaining aspect ratio
                        if max_width / max_height > aspect_ratio:
                            # Height is the limiting factor
                            actual_height = max_height
                            actual_width = Inches(actual_height / Inches(1) * aspect_ratio)
                        else:
                            # Width is the limiting factor
                            actual_width = max_width
                            actual_height = Inches(actual_width / Inches(1) / aspect_ratio)

                        # Center the image within the max bounds
                        left = max_left + (max_width - actual_width) / 2
                        top = max_top + (max_height - actual_height) / 2

                        slide.shapes.add_picture(
                            BytesIO(image_bytes),
                            left,
                            top,
                            width=actual_width,
                            height=actual_height,
                        )
                        logger.info(
                            "Successfully added image to slide for item %s (%s)",
                            item.id,
                            item.title,
                        )
                    except (httpx.HTTPError, Exception) as exc:
                        logger.warning(
                            "Unable to fetch or process slide image for auction item %s (%s): %s",
                            item.id,
                            item.title,
                            str(exc),
                        )
                else:
                    logger.warning(
                        "No primary_image_url for auction item %s (%s) - slide will have no image",
                        item.id,
                        item.title,
                    )

                # Add semi-transparent background for ON_IMAGE layout
                if layout_key == "on_image":
                    # Create a rounded rectangle covering the text area
                    bg_left = Inches(0.6)
                    bg_top = Inches(5.6)
                    bg_width = Inches(12.133)
                    bg_height = Inches(1.7)

                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        bg_left,
                        bg_top,
                        bg_width,
                        bg_height,
                    )

                    # Set semi-transparent black fill (70% opacity)
                    bg_shape.fill.solid()
                    bg_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)  # type: ignore[no-untyped-call]
                    bg_shape.fill.transparency = 0.3  # 30% transparent = 70% opacity

                    # Remove border
                    bg_shape.line.fill.background()

                # Add title textbox
                title_left, title_top, title_width, title_height = config["title"]
                title_box = slide.shapes.add_textbox(
                    title_left, title_top, title_width, title_height
                )
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_para = title_frame.paragraphs[0]
                title_para.text = f"#{item.bid_number or '—'}  {item.title}"
                title_para.font.size = Pt(24)
                title_para.font.bold = True

                # Set white text for ON_IMAGE layout
                if layout_key == "on_image":
                    title_para.font.color.rgb = RGBColor(255, 255, 255)  # type: ignore[no-untyped-call]

                # Add text content textbox
                text_left, text_top, text_width, text_height = config["text"]
                text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.text = body_text
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    # Set white text for ON_IMAGE layout
                    if layout_key == "on_image":
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # type: ignore[no-untyped-call]

        output = BytesIO()
        presentation.save(output)
        output.seek(0)
        return output.read()
